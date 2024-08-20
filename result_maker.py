import cv2
import os
import json
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import json
from PIL import Image
import numpy as np
import matplotlib.pyplot as plt
from matplotlib import cm
import datetime
import tkinter as tk
from tkinter import filedialog, ttk
from tkinter import font as tkfont
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")


# パラメータリストファイルの読み込み
with open("parameters_list.json", "r") as f:
    params_list = json.load(f)

# 各パラメータセットに対して処理を実行
for params in params_list:
    # 入力動画フォルダと出力画像フォルダの設定
    input_folder = params.get("input_folder", "./input_video")
    output_folder = params.get("output_folder", "./output_images")

    # 出力フォルダが存在しない場合は作成
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # frame_intervalを直接読み込む
    frame_interval = params.get(
        "frame_interval", 10800
    )  # デフォルト値は10800（360秒 * 30fps）

    # 動画ファイルの処理
    for video_file in os.listdir(input_folder):
        if video_file.endswith((".mp4", ".avi", ".mov")):  # 対応する動画形式
            video_path = os.path.join(input_folder, video_file)
            cap = cv2.VideoCapture(video_path)

            frame_count = 0
            image_count = 0
            max_images = params.get("images_per_video", 8)

            while True:
                ret, frame = cap.read()
                if not ret:
                    break

                if frame_count % frame_interval == 0 or (
                    params.get("include_first_frame", True) and frame_count == 0
                ):
                    output_path = os.path.join(
                        output_folder, f"{video_file[:-4]}_{frame_count:06d}.jpg"
                    )
                    cv2.imwrite(output_path, frame)
                    image_count += 1

                    if image_count >= max_images:
                        break

                frame_count += 1

            cap.release()

print("すべてのパラメータセットに対する処理が完了しました。")


# 定数の定義
SLIDE_WIDTH = Cm(21)
SLIDE_HEIGHT = Cm(29.7)
TABLE_HEIGHT = Cm(11)  # 表の高さを短くする


def load_parameters(file_path):
    with open(file_path, "r") as f:
        return json.load(f)


def create_colorbar(min_value, max_value):
    fig, ax = plt.subplots(figsize=(0.5, 5))  # 幅を0.5に変更してカラーバーを細くする
    gradient = np.linspace(min_value, max_value, 256).reshape(256, 1)
    ax.imshow(
        gradient, aspect="auto", cmap="jet_r", extent=[0, 1, min_value, max_value]
    )
    ax.yaxis.set_label_position("right")
    ax.xaxis.set_visible(False)
    ax.yaxis.tick_right()

    # フォントサイズを大きくする
    ax.tick_params(axis="y", labelsize=20)  # ここでフォントサイズを調整

    # 上限値と下限値を必ず表示する
    ax.set_yticks([min_value, max_value])
    ax.set_yticklabels([f"{min_value:.1f}", f"{max_value:.1f}"])

    plt.tight_layout()

    temp_file = "temp_colorbar.png"
    plt.savefig(temp_file, dpi=300, bbox_inches="tight")
    plt.close()

    return temp_file


def group_images_by_video(image_folder, include_first_frame):
    image_files = [
        f for f in os.listdir(image_folder) if f.endswith((".jpg", ".png", ".jpeg"))
    ]
    image_groups = {}
    for image_file in image_files:
        video_name = image_file.split("_")[0]
        frame_number = int(image_file.split("_")[-1].split(".")[0])
        if not include_first_frame and frame_number == 0:
            continue
        if video_name not in image_groups:
            image_groups[video_name] = []
        image_groups[video_name].append(image_file)
    return image_groups


def add_grid_lines(
    slide, table_left, table_width, title_top, cell_width, cell_height, rows, cols
):
    for row in range(rows + 1):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left,
            title_top + row * cell_height,
            table_width,
            0,
        )
        line.fill.background()
        line.line.color.rgb = RGBColor(0, 0, 0)

    for col in range(cols + 1):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left + col * cell_width,
            title_top,
            0,
            rows * cell_height,
        )
        line.fill.background()
        line.line.color.rgb = RGBColor(0, 0, 0)


def calculate_table_height(cell_width, rows, image_aspect_ratio):
    cell_height = cell_width / image_aspect_ratio
    frame_text_height = Cm(0.7)  # フレーム数の文字の高さ
    return (cell_height + frame_text_height) * rows


def format_time(seconds):
    return str(datetime.timedelta(seconds=seconds)).split(".")[0]


def add_image_to_slide(
    slide, img_path, left, top, cell_width, cell_height, seconds_per_frame
):
    with Image.open(img_path) as img:
        img_width, img_height = img.size

    aspect_ratio = img_width / img_height
    image_height = cell_height - Cm(0.7)  # フレーム数の文字の高さを引く
    image_width = image_height * aspect_ratio

    if image_width > cell_width:
        image_width = cell_width
        image_height = image_width / aspect_ratio

    image_left = left + (cell_width - image_width) / 2
    image_top = top

    slide.shapes.add_picture(
        img_path, image_left, image_top, width=image_width, height=image_height
    )

    frame_number = int(os.path.basename(img_path).split("_")[-1].split(".")[0])
    elapsed_time = format_time(frame_number * seconds_per_frame)
    txBox = slide.shapes.add_textbox(left, top + image_height, cell_width, Cm(0.7))
    tf = txBox.text_frame
    tf.text = f"{elapsed_time.split(':')[0]}:{elapsed_time.split(':')[1]}, frame:{frame_number}"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)  # フォントサイズを小さくして2行に収める


def create_presentation(params_list):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_slide_layout = prs.slide_layouts[6]

    # グローバルパラメータの設定（最初の要素をデフォルトとして使用）
    global_params = params_list[0] if params_list else {}

    # デフォルト値の設定
    default_params = {
        "show_colorbar": True,
        "images_per_video": 8,
        "rows": 2,
        "cols": 4,
        "include_first_frame": True,
        "seconds_per_frame": 360,
        "min_threshold": 0,
        "max_threshold": 255,
        "output_folder": "",
    }

    # グローバルパラメータにデフォルト値を適用
    for key, value in default_params.items():
        if key not in global_params:
            global_params[key] = value

    image_groups = group_images_by_video(
        global_params["output_folder"], global_params["include_first_frame"]
    )
    video_names = list(image_groups.keys())

    for i, params in enumerate(params_list):
        if i % 2 == 0:
            slide = prs.slides.add_slide(blank_slide_layout)

        video_index = i % 2

        # 現在のパラメータにグローバルパラメータを適用し、その後個別のパラメータで上書き
        current_params = global_params.copy()
        current_params.update(params)

        table_left = Cm(0.5)
        available_height = SLIDE_HEIGHT / 2 - Cm(2)

        if current_params["show_colorbar"]:
            colorbar_file = create_colorbar(
                current_params["min_threshold"], current_params["max_threshold"]
            )
            colorbar_left, colorbar_top = Cm(0.5), video_index * (
                SLIDE_HEIGHT / 2
            ) + Cm(1)
            colorbar_width, colorbar_height = (
                Cm(1.5),
                available_height / 2,
            )  # 高さを半分に変更
            slide.shapes.add_picture(
                colorbar_file,
                colorbar_left,
                colorbar_top,
                width=colorbar_width,
                height=colorbar_height,
            )
            os.remove(colorbar_file)
            table_left = colorbar_left + colorbar_width + Cm(0.5)

        table_width = SLIDE_WIDTH - table_left - Cm(0.5)
        cell_width = table_width / current_params["cols"]

        # 画像のアスペクト比を取得（最初の画像を使用）
        if i < len(video_names):
            video_name = video_names[i]
            first_image = os.path.join(
                current_params["output_folder"], image_groups[video_name][0]
            )
            with Image.open(first_image) as img:
                image_aspect_ratio = img.width / img.height
        else:
            # ビデオが存在しない場合のデフォルト値
            image_aspect_ratio = 16 / 9  # 一般的な動画のアスペクト比

        table_height = calculate_table_height(
            cell_width, current_params["rows"], image_aspect_ratio
        )
        cell_height = table_height / current_params["rows"]

        title_top = video_index * (SLIDE_HEIGHT / 2)
        title = slide.shapes.add_textbox(Cm(1), title_top, SLIDE_WIDTH - Cm(2), Cm(1))
        if i < len(video_names):
            title.text = video_names[i]
        else:
            title.text = f"ビデオ {i+1} (データなし)"
        title.text_frame.paragraphs[0].font.size = Pt(18)
        title.text_frame.paragraphs[0].font.bold = True

        add_grid_lines(
            slide,
            table_left,
            table_width,
            title_top + Cm(1),
            cell_width,
            cell_height,
            current_params["rows"],
            current_params["cols"],
        )

        if i < len(video_names):
            group_images = sorted(
                image_groups[video_name],
                key=lambda x: int(x.split("_")[-1].split(".")[0]),
            )
            for j, img_file in enumerate(
                group_images[: current_params["images_per_video"]]
            ):
                img_path = os.path.join(current_params["output_folder"], img_file)
                row, col = j // current_params["cols"], j % current_params["cols"]
                left = table_left + col * cell_width
                top = title_top + Cm(1) + row * cell_height
                add_image_to_slide(
                    slide,
                    img_path,
                    left,
                    top,
                    cell_width,
                    cell_height,
                    current_params["seconds_per_frame"],
                )

    prs.save("image_summary_a4_two_videos_improved.pptx")


def create_gui():
    root = tk.Tk()
    root.title("Create Presentation")

    # フォントを指定してUTF-8エンコーディングを使用
    default_font = tkfont.nametofont("TkDefaultFont")
    default_font.configure(size=10, family="MS Gothic")
    root.option_add("*Font", default_font)

    def browse_folder():
        folder_path = filedialog.askdirectory()
        output_folder_entry.delete(0, tk.END)
        output_folder_entry.insert(0, folder_path)

    def on_execute():
        params = {
            "output_folder": output_folder_entry.get(),
            "images_per_video": int(images_per_video_entry.get()),
            "rows": int(rows_entry.get()),
            "cols": int(cols_entry.get()),
            "include_first_frame": include_first_frame_var.get(),
            "seconds_per_frame": int(seconds_per_frame_entry.get()),
            "min_threshold": float(min_threshold_entry.get()),
            "max_threshold": float(max_threshold_entry.get()),
            "show_colorbar": show_colorbar_var.get(),
        }
        root.quit()
        root.destroy()
        create_presentation([params])
        print("Presentation has been created.")

    root = tk.Tk()
    root.title("Create Presentation")

    frame = ttk.Frame(root, padding="10")
    frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))

    ttk.Label(frame, text="Output Folder").grid(column=0, row=0, sticky=tk.W)
    output_folder_entry = ttk.Entry(frame, width=50)
    output_folder_entry.grid(column=1, row=0, sticky=(tk.W, tk.E))
    ttk.Button(frame, text="Browse", command=browse_folder).grid(
        column=2, row=0, sticky=tk.W
    )

    ttk.Label(frame, text="Images per Video").grid(column=0, row=1, sticky=tk.W)
    images_per_video_entry = ttk.Entry(frame, width=10)
    images_per_video_entry.insert(0, "8")
    images_per_video_entry.grid(column=1, row=1, sticky=tk.W)

    ttk.Label(frame, text="Rows").grid(column=0, row=2, sticky=tk.W)
    rows_entry = ttk.Entry(frame, width=10)
    rows_entry.insert(0, "2")
    rows_entry.grid(column=1, row=2, sticky=tk.W)

    ttk.Label(frame, text="Columns").grid(column=0, row=3, sticky=tk.W)
    cols_entry = ttk.Entry(frame, width=10)
    cols_entry.insert(0, "4")
    cols_entry.grid(column=1, row=3, sticky=tk.W)

    include_first_frame_var = tk.BooleanVar(value=True)
    ttk.Checkbutton(
        frame, text="Include First Frame", variable=include_first_frame_var
    ).grid(column=0, row=4, columnspan=2, sticky=tk.W)

    ttk.Label(frame, text="Seconds per Frame").grid(column=0, row=5, sticky=tk.W)
    seconds_per_frame_entry = ttk.Entry(frame, width=10)
    seconds_per_frame_entry.insert(0, "360")
    seconds_per_frame_entry.grid(column=1, row=5, sticky=tk.W)

    ttk.Label(frame, text="Min Threshold").grid(column=0, row=6, sticky=tk.W)
    min_threshold_entry = ttk.Entry(frame, width=10)
    min_threshold_entry.insert(0, "0")
    min_threshold_entry.grid(column=1, row=6, sticky=tk.W)

    ttk.Label(frame, text="Max Threshold").grid(column=0, row=7, sticky=tk.W)
    max_threshold_entry = ttk.Entry(frame, width=10)
    max_threshold_entry.insert(0, "255")
    max_threshold_entry.grid(column=1, row=7, sticky=tk.W)

    show_colorbar_var = tk.BooleanVar(value=True)
    ttk.Checkbutton(frame, text="Show Colorbar", variable=show_colorbar_var).grid(
        column=0, row=8, columnspan=2, sticky=tk.W
    )

    ttk.Button(frame, text="Execute", command=on_execute).grid(
        column=0, row=9, sticky=tk.W
    )
    ttk.Button(frame, text="Cancel", command=root.quit).grid(
        column=1, row=9, sticky=tk.E
    )

    root.mainloop()


if __name__ == "__main__":
    create_gui()
