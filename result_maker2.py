import cv2
import os
import json
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
import numpy as np
import matplotlib.pyplot as plt
from matplotlib import cm
import datetime
import tkinter as tk
from tkinter import filedialog, ttk
from tkinter import font as tkfont
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

# 定数の定義
SLIDE_WIDTH = Cm(21)
SLIDE_HEIGHT = Cm(29.7)
TABLE_HEIGHT = Cm(11)  # 表の高さを短くする


def select_video_directory():
    directory = filedialog.askdirectory()
    video_directory_entry.delete(0, tk.END)
    video_directory_entry.insert(0, directory)
    update_video_list()


def update_video_list():
    video_list.delete(0, tk.END)
    directory = video_directory_entry.get()
    for file in os.listdir(directory):
        if file.endswith((".mp4", ".avi", ".mov")):
            video_list.insert(tk.END, file)


def extract_frames():
    directory = video_directory_entry.get()
    output_directory = "./output_images"
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)

    extracted_frames = {}

    for video_file in video_list.get(0, tk.END):
        interval = int(frame_interval_entry.get())
        video_path = os.path.join(directory, video_file)
        cap = cv2.VideoCapture(video_path)

        frame_count = 0
        image_count = 0

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % interval == 0:
                output_path = os.path.join(
                    output_directory, f"{video_file[:-4]}_{frame_count:06d}.jpg"
                )
                cv2.imwrite(output_path, frame)
                image_count += 1

            frame_count += 1

        cap.release()
        extracted_frames[video_file] = image_count
        print(f"{video_file}: Extracted {image_count} images.")

    show_slide_parameters_gui(extracted_frames)


def show_slide_parameters_gui(extracted_frames):
    params_window = tk.Toplevel(root)
    params_window.title("Slide Parameters Settings")

    params_list = []

    def add_param_form(video_name, image_count):
        param_frame = ttk.Frame(params_window, padding="10")
        param_frame.pack(fill=tk.X, expand=True)

        ttk.Label(
            param_frame, text=f"Video: {video_name} (Extracted images: {image_count})"
        ).grid(column=0, row=0, columnspan=3, sticky=tk.W)

        ttk.Label(param_frame, text="Images per Video").grid(
            column=0, row=1, sticky=tk.W
        )
        param_frame.images_per_video_entry = ttk.Entry(param_frame, width=10)
        param_frame.images_per_video_entry.insert(0, str(min(8, image_count)))
        param_frame.images_per_video_entry.grid(column=1, row=1, sticky=tk.W)

        ttk.Label(param_frame, text="Rows").grid(column=0, row=2, sticky=tk.W)
        param_frame.rows_entry = ttk.Entry(param_frame, width=10)
        param_frame.rows_entry.insert(0, "2")
        param_frame.rows_entry.grid(column=1, row=2, sticky=tk.W)

        ttk.Label(param_frame, text="Columns").grid(column=0, row=3, sticky=tk.W)
        param_frame.cols_entry = ttk.Entry(param_frame, width=10)
        param_frame.cols_entry.insert(0, "4")
        param_frame.cols_entry.grid(column=1, row=3, sticky=tk.W)

        param_frame.include_first_frame_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(
            param_frame,
            text="Include First Frame",
            variable=param_frame.include_first_frame_var,
        ).grid(column=0, row=4, columnspan=2, sticky=tk.W)

        ttk.Label(param_frame, text="Seconds per Frame").grid(
            column=0, row=5, sticky=tk.W
        )
        param_frame.seconds_per_frame_entry = ttk.Entry(param_frame, width=10)
        param_frame.seconds_per_frame_entry.insert(0, "360")
        param_frame.seconds_per_frame_entry.grid(column=1, row=5, sticky=tk.W)

        ttk.Label(param_frame, text="Min Threshold").grid(column=0, row=6, sticky=tk.W)
        param_frame.min_threshold_entry = ttk.Entry(param_frame, width=10)
        param_frame.min_threshold_entry.insert(0, "0")
        param_frame.min_threshold_entry.grid(column=1, row=6, sticky=tk.W)

        ttk.Label(param_frame, text="Max Threshold").grid(column=0, row=7, sticky=tk.W)
        param_frame.max_threshold_entry = ttk.Entry(param_frame, width=10)
        param_frame.max_threshold_entry.insert(0, "255")
        param_frame.max_threshold_entry.grid(column=1, row=7, sticky=tk.W)

        param_frame.show_colorbar_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(
            param_frame, text="Show Colorbar", variable=param_frame.show_colorbar_var
        ).grid(column=0, row=8, columnspan=2, sticky=tk.W)

        return param_frame

    for video_name, image_count in extracted_frames.items():
        params_list.append(add_param_form(video_name, image_count))

    def on_execute():
        final_params = []
        for frame in params_list:
            params = {
                "output_folder": "./output_images",
                "images_per_video": int(frame.images_per_video_entry.get()),
                "rows": int(frame.rows_entry.get()),
                "cols": int(frame.cols_entry.get()),
                "include_first_frame": frame.include_first_frame_var.get(),
                "seconds_per_frame": int(frame.seconds_per_frame_entry.get()),
                "min_threshold": float(frame.min_threshold_entry.get()),
                "max_threshold": float(frame.max_threshold_entry.get()),
                "show_colorbar": frame.show_colorbar_var.get(),
            }
            final_params.append(params)

        params_window.destroy()
        create_presentation(final_params)

    ttk.Button(params_window, text="Execute", command=on_execute).pack()


def create_colorbar(min_value, max_value):
    fig, ax = plt.subplots(figsize=(0.5, 5))  # 幅を0.5に変更してカラーバーを細くする
    gradient = np.linspace(min_value, max_value, 256).reshape(256, 1)
    ax.imshow(
        gradient, aspect="auto", cmap="jet_r", extent=[0, 1, min_value, max_value]
    )
    ax.yaxis.set_label_position("right")
    ax.xaxis.set_visible(False)
    ax.yaxis.tick_right()

    # フォントサイズを大きくする
    ax.tick_params(axis="y", labelsize=20)  # ここでフォントサイズを調整

    # 上限値と下限値を必ず表示する
    ax.set_yticks([min_value, max_value])
    ax.set_yticklabels([f"{min_value:.1f}", f"{max_value:.1f}"])

    temp_file = "temp_colorbar.png"
    plt.savefig(temp_file, dpi=300, bbox_inches="tight")
    plt.close()

    return temp_file


def group_images_by_video(image_folder, include_first_frame):
    image_files = [
        f for f in os.listdir(image_folder) if f.endswith((".jpg", ".png", ".jpeg"))
    ]
    image_groups = {}
    for image_file in image_files:
        video_name = image_file.split("_")[0]
        frame_number = int(image_file.split("_")[-1].split(".")[0])
        if not include_first_frame and frame_number == 0:
            continue
        if video_name not in image_groups:
            image_groups[video_name] = []
        image_groups[video_name].append(image_file)
    return image_groups


def add_grid_lines(
    slide, table_left, table_width, title_top, cell_width, cell_height, rows, cols
):
    for row in range(rows + 1):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left,
            title_top + row * cell_height,
            table_width,
            0,
        )
        line.fill.background()
        line.line.color.rgb = RGBColor(0, 0, 0)

    for col in range(cols + 1):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left + col * cell_width,
            title_top,
            0,
            rows * cell_height,
        )
        line.fill.background()
        line.line.color.rgb = RGBColor(0, 0, 0)


def calculate_table_height(cell_width, rows, image_aspect_ratio):
    cell_height = cell_width / image_aspect_ratio
    frame_text_height = Cm(0.7)  # フレーム数の文字の高さ
    return (cell_height + frame_text_height) * rows


def format_time(seconds):
    return str(datetime.timedelta(seconds=seconds)).split(".")[0]


def add_image_to_slide(
    slide, img_path, left, top, cell_width, cell_height, seconds_per_frame
):
    with Image.open(img_path) as img:
        img_width, img_height = img.size

    aspect_ratio = img_width / img_height
    frame_text_height = Cm(0.7)  # フレーム数の文字の高さ
    available_height = cell_height - frame_text_height

    # セル内の余白を設定（上下左右均等に）
    margin = min(cell_width, available_height) * 0.05

    # 画像の最大サイズを計算
    max_width = cell_width - 2 * margin
    max_height = available_height - 2 * margin

    # アスペクト比を維持しながら、画像サイズを調整
    if aspect_ratio > max_width / max_height:
        image_width = max_width
        image_height = image_width / aspect_ratio
    else:
        image_height = max_height
        image_width = image_height * aspect_ratio

    # 画像の配置位置を計算（セル内で中央揃え）
    image_left = left + (cell_width - image_width) / 2
    image_top = top + (available_height - image_height) / 2

    slide.shapes.add_picture(
        img_path, image_left, image_top, width=image_width, height=image_height
    )

    frame_number = int(os.path.basename(img_path).split("_")[-1].split(".")[0])
    elapsed_time = format_time(frame_number * seconds_per_frame)
    txBox = slide.shapes.add_textbox(
        left, top + available_height, cell_width, frame_text_height
    )
    tf = txBox.text_frame
    tf.text = f"{elapsed_time.split(':')[0]}:{elapsed_time.split(':')[1]}, frame:{frame_number}"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)  # フォントサイズを小さくして2行に収める


def create_presentation(params_list):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_slide_layout = prs.slide_layouts[6]

    image_groups = group_images_by_video(
        "./output_images", True
    )  # すべての画像を含める
    video_names = list(image_groups.keys())

    for i, params in enumerate(params_list):
        if i % 2 == 0:
            slide = prs.slides.add_slide(blank_slide_layout)

        video_index = i % 2

        table_left = Cm(0.5)
        available_height = SLIDE_HEIGHT / 2 - Cm(
            2
        )  # スライドの半分から余白を引いた高さ

        if params["show_colorbar"]:
            # ... カラーバーの処理 ...
            table_left = colorbar_left + colorbar_width + Cm(0.5)

        max_table_width = SLIDE_WIDTH - table_left - Cm(0.5)
        initial_cell_width = max_table_width / params["cols"]

        if i < len(video_names):
            video_name = video_names[i]
            first_image = os.path.join("./output_images", image_groups[video_name][0])
            with Image.open(first_image) as img:
                image_aspect_ratio = img.width / img.height
        else:
            image_aspect_ratio = 16 / 9

        # 表の高さを計算し、利用可能な高さを超えないように調整
        table_height = calculate_table_height(
            initial_cell_width, params["rows"], image_aspect_ratio
        )
        if table_height > available_height:
            scale_factor = available_height / table_height
            table_height = available_height
            cell_height = table_height / params["rows"]
            cell_width = initial_cell_width * scale_factor
            table_width = cell_width * params["cols"]
        else:
            cell_height = table_height / params["rows"]
            cell_width = initial_cell_width
            table_width = max_table_width

        title_top = video_index * (SLIDE_HEIGHT / 2)
        title = slide.shapes.add_textbox(Cm(1), title_top, SLIDE_WIDTH - Cm(2), Cm(1))
        if i < len(video_names):
            title.text = video_names[i]
        else:
            title.text = f"ビデオ {i+1} (データなし)"
        title.text_frame.paragraphs[0].font.size = Pt(18)
        title.text_frame.paragraphs[0].font.bold = True

        add_grid_lines(
            slide,
            table_left,
            table_width,
            title_top + Cm(1),
            cell_width,
            cell_height,
            params["rows"],
            params["cols"],
        )

        if i < len(video_names):
            group_images = sorted(
                image_groups[video_name],
                key=lambda x: int(x.split("_")[-1].split(".")[0]),
            )
            if not params["include_first_frame"]:
                group_images = group_images[1:]
            for j, img_file in enumerate(group_images[: params["images_per_video"]]):
                if j >= params["rows"] * params["cols"]:
                    break
                img_path = os.path.join("./output_images", img_file)
                row, col = j // params["cols"], j % params["cols"]
                left = table_left + col * cell_width
                top = title_top + Cm(1) + row * cell_height
                add_image_to_slide(
                    slide,
                    img_path,
                    left,
                    top,
                    cell_width,
                    cell_height,
                    params["seconds_per_frame"],
                )
    現在の日時 = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    prs.save(f"image_summary_{現在の日時}.pptx")
    print("プレゼンテーションが作成されました。")
    root.quit()


def create_custom_style():
    style = ttk.Style()
    style.configure("TLabel", font=("Helvetica", 15))
    style.configure("TButton", font=("Helvetica", 15))
    style.configure("TEntry", font=("Helvetica", 15))
    return style


# メインのGUI
root = tk.Tk()
root.title("Video Frame Extraction")

# カスタムスタイルを適用
custom_style = create_custom_style()

# 以下の行を変更
ttk.Label(root, text="Video Directory:", style="TLabel").grid(
    column=0, row=0, sticky=tk.W
)
video_directory_entry = ttk.Entry(root, width=50, font=("Helvetica", 15))
video_directory_entry.grid(column=1, row=0)
ttk.Button(root, text="Browse", command=select_video_directory, style="TButton").grid(
    column=2, row=0
)

ttk.Label(root, text="Frame Interval:", style="TLabel").grid(
    column=0, row=1, sticky=tk.W
)
frame_interval_entry = ttk.Entry(root, width=10, font=("Helvetica", 15))
frame_interval_entry.grid(column=1, row=1, sticky=tk.W)
frame_interval_entry.insert(0, "30")  # Default value

video_list = tk.Listbox(root, width=50, height=10, font=("Helvetica", 15))
video_list.grid(column=0, row=2, columnspan=3)

ttk.Button(root, text="Extract Frames", command=extract_frames, style="TButton").grid(
    column=1, row=3
)


root.mainloop()
